import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
SLIDE_WIDTH = Inches(22)
SLIDE_HEIGHT = Inches(13)
FONT_NAME = "Arial"

# --- CONSTANTS ---
# Orientation constant (Manual definition to avoid ImportError)
MSO_TEXT_ORIENTATION_UPWARD = 2 

# --- COLORS ---
# Header Backgrounds (Darker)
C_BLUE_DARK = RGBColor(0, 32, 96)       # Deep Blue (Security, Data Headers)
C_PURPLE_DARK = RGBColor(89, 38, 128)   # Darker Purple
C_GREY_DARK = RGBColor(64, 64, 64)      # Dark Grey (Client Apps, Infra)
C_RED_DARK = RGBColor(192, 0, 0)        # Darker Red for OpenShift Border

# Product Box Borders (Thinner, slightly lighter than header)
C_BLUE_BORDER = RGBColor(0, 51, 153)
C_PURPLE_BORDER = RGBColor(112, 48, 160)
C_GREY_BORDER = RGBColor(89, 89, 89)

# Standard Colors
C_WHITE = RGBColor(255, 255, 255)
C_BLACK = RGBColor(0, 0, 0)
C_GREY_LIGHT = RGBColor(242, 242, 242)  # Backgrounds
C_BTN_BLUE = RGBColor(68, 114, 196)     # Button Blue

# Legend Colors
C_LEG_ENTITLED = RGBColor(189, 215, 238)
C_LEG_OPP = RGBColor(169, 209, 142)
C_LEG_EXPLORE = RGBColor(255, 217, 102)
C_LEG_RISK = RGBColor(244, 176, 132)
C_ELA_STAR = RGBColor(217, 217, 217)

# --- HELPER FUNCTIONS ---

def create_box(slide, x, y, w, h, text, bg_color, font_color=C_BLACK, 
               bold=False, font_size=8, outline_color=None, outline_width=Pt(0.75),
               align=PP_ALIGN.CENTER, vertical_text=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    
    if outline_color:
        shape.line.color.rgb = outline_color
        shape.line.width = outline_width
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    # TIGHT MARGINS to match the "pitch" and fit long text
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if vertical_text:
        tf.orientation = MSO_TEXT_ORIENTATION_UPWARD
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.alignment = align
    return shape

def build_product_grid(slide, x, y, w, h, products, cols=1, border_color=C_BLUE_BORDER):
    gap = Inches(0.05)
    if not products: return
    rows = (len(products) + cols - 1) // cols
    box_w = (w - (gap * (cols - 1))) / cols
    box_h = (h - (gap * (rows - 1))) / rows
    
    # Cap height for aesthetics so boxes aren't huge if list is short
    if box_h > Inches(0.45): box_h = Inches(0.45)

    for i, prod in enumerate(products):
        r = i // cols
        c = i % cols
        bx = x + (c * (box_w + gap))
        by = y + (r * (box_h + gap))
        # Product boxes: White background, colored border, thinner line
        # Font Size 8pt fits the density best
        create_box(slide, bx, by, box_w, box_h, prod, C_WHITE, outline_color=border_color, outline_width=Pt(1.0), font_size=8)

# --- MAIN SCRIPT ---

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. HEADER AREA
create_box(slide, Inches(0.2), Inches(0.2), Inches(3), Inches(0.6), "IBM Technology", C_WHITE, font_size=24, bold=True, align=PP_ALIGN.LEFT)

# Legends
leg_w = Inches(1.8)
leg_x = SLIDE_WIDTH - (leg_w * 4) - Inches(0.5)
legends = [("Entitled", C_LEG_ENTITLED), ("Opportunity", C_LEG_OPP), ("Explore", C_LEG_EXPLORE), ("No Interest/At Risk", C_LEG_RISK)]
for txt, col in legends:
    create_box(slide, leg_x, Inches(0.35), leg_w, Inches(0.35), txt, col, font_size=11, bold=False)
    leg_x += leg_w + Inches(0.1)

create_box(slide, SLIDE_WIDTH - Inches(2), Inches(0.8), Inches(1.5), Inches(0.3), "★ ELA Product", C_ELA_STAR, font_size=10, bold=True)

# --- LAYOUT DEFINITIONS ---
MARGIN_X = Inches(0.2)
TOP_Y = Inches(1.2)
BOTTOM_FOOTER_H = Inches(0.6)
MAIN_H = SLIDE_HEIGHT - TOP_Y - BOTTOM_FOOTER_H - Inches(0.2)

# Column Widths
CE_W = Inches(0.4)       # Far Left: Client Engineering
SEC_W = Inches(2.3)      # Security Column
RIGHT_W = Inches(2.3)    # Right Column (IT Auto)
GAP = Inches(0.1)

# Center Width
CENTER_W = SLIDE_WIDTH - (MARGIN_X * 2) - CE_W - SEC_W - RIGHT_W - (GAP * 3)

# X Coordinates
X_CE = MARGIN_X
X_SEC = X_CE + CE_W + GAP
X_CENTER = X_SEC + SEC_W + GAP
X_RIGHT = X_CENTER + CENTER_W + GAP

# --- 2. FAR LEFT: CLIENT ENGINEERING ---
# Increased font size to 10pt for visibility
ce_box = create_box(slide, X_CE, TOP_Y + Inches(0.3), CE_W, MAIN_H, "IBM Client Engineering (CE)", C_WHITE, outline_color=C_BLACK, outline_width=Pt(1.5), bold=True, font_size=11, vertical_text=True)

# --- 3. LEFT COLUMN: SECURITY ---
# Security Header
create_box(slide, X_SEC, TOP_Y, SEC_W, Inches(0.3), "Security", C_WHITE, font_size=11, bold=True)

# Data Security Block
DS_H = Inches(4.2)
# Header: 9pt Bold
create_box(slide, X_SEC, TOP_Y + Inches(0.35), SEC_W, Inches(0.35), "Data Security", C_BLUE_DARK, font_color=C_WHITE, font_size=9, bold=True)
ds_prods = ["Guardium Data Encryption", "Guardium Data Protection", "Guardium Data Security Center", "Guardium Discover and Classify", "Guardium Key Lifecycle Management"]
build_product_grid(slide, X_SEC, TOP_Y + Inches(0.75), SEC_W, DS_H - Inches(0.4), ds_prods, cols=1, border_color=C_BLUE_BORDER)

# Identity Block
ID_Y = TOP_Y + Inches(0.35) + DS_H + GAP
ID_H = Inches(4.5)
create_box(slide, X_SEC, ID_Y, SEC_W, Inches(0.35), "Identity & Access Mgmt", C_PURPLE_DARK, font_color=C_WHITE, font_size=9, bold=True)
id_prods = ["HashiCorp Boundary", "HashiCorp Consul", "HashiCorp Vault", "ILMT", "Security Verify (IAM)", "Security MaaS 360", "Trusteer (Anti-fraud)"]
build_product_grid(slide, X_SEC, ID_Y + Inches(0.4), SEC_W, ID_H - Inches(0.4), id_prods, cols=1, border_color=C_PURPLE_BORDER)


# --- 4. CENTER BLOCK ---
# A. Client Apps (Top)
CA_H = Inches(1.4)
create_box(slide, X_CENTER, TOP_Y, CENTER_W, Inches(0.3), "Client Applications", C_GREY_DARK, font_color=C_WHITE, font_size=11, bold=True)
ca_prods = ["ERP", "CRM", "B2B", "B2C", "B2E", "Omnichannel", "CRM (on-prem)", "IA", "Fraud", "Credit", "PCP", "Supply Chain", "Engineering / Network", "Portal / Mobile / APP", "Payment Instantaneous", "Customer Service"]
build_product_grid(slide, X_CENTER, TOP_Y + Inches(0.35), CENTER_W, CA_H - Inches(0.35), ca_prods, cols=8, border_color=C_GREY_BORDER)

# B. 6 Pillars (Middle)
P6_Y = TOP_Y + CA_H + GAP
P6_H = Inches(5.0)

# Data Structure: Title, Header Color, Border Color, Products, Columns
pillars = [
    ("AI Assistants", C_BLUE_DARK, C_BLUE_BORDER, ["Automation", "Blueworks Live", "Business Analytics", "Business Automation", "CP4BA", "Cognos Analytics", "Decision Mgmt", "Planning Analytics", "Process Mining", "RPA", "SPSS Modeler", "watsonx Assistants", "watsonx BI Assistant", "watsonx Code Assistant", "watsonx Orchestrate", "Workflow Automation"], 2),
    ("AI/MLOps", C_BLUE_DARK, C_BLUE_BORDER, ["CP4D", "OpenPages", "Orchestrate (SaaS)", "WCA Ansible & Java", "WCAz", "watsonx.ai", "watsonx.governance"], 1),
    ("Databases", C_BLUE_DARK, C_BLUE_BORDER, ["CM8", "CMOD", "CP4D", "Capture", "Cloudera", "Content", "DB2", "Database Eco", "FileNet", "Hadoop", "Informix", "Netezza", "watsonx.data", "watsonx.ai (SaaS)"], 2),
    ("Data Intelligence", C_BLUE_DARK, C_BLUE_BORDER, ["CP4D", "Data Product Hub", "Decision Optimization", "Knowledge Catalog", "Manta Data Lineage", "Optim & Master Data Mgmt", "SPSS Stats"], 1),
    ("Data Integration", C_BLUE_DARK, C_BLUE_BORDER, ["CP4D", "Data Fabric", "Data Integration", "DataStage", "Databand", "Replication", "StreamSets"], 1),
    ("Asset Lifecycle Management", C_PURPLE_DARK, C_PURPLE_BORDER, ["EI", "Envizi", "HashiCorp Terraform", "Maximo", "Sterling Order & Inventory Mgmt", "Supply Chain", "TRIRIGA"], 1)
]
p_w = (CENTER_W - (GAP * 5)) / 6
for i, (title, h_col, b_col, prods, col_count) in enumerate(pillars):
    px = X_CENTER + (i * (p_w + GAP))
    create_box(slide, px, P6_Y, p_w, Inches(0.35), title, h_col, font_color=C_WHITE, font_size=9, bold=True)
    build_product_grid(slide, px, P6_Y + Inches(0.4), p_w, P6_H - Inches(0.4), prods, cols=col_count, border_color=b_col)

# C. App Dev & Integration (Bottom of Center)
AD_Y = P6_Y + P6_H + GAP
AD_H = Inches(1.8)
ad_w = (CENTER_W - GAP) / 2
# App Dev
create_box(slide, X_CENTER, AD_Y, ad_w, Inches(0.3), "Application Development", C_PURPLE_DARK, font_color=C_WHITE, font_size=11, bold=True)
ad_prods = ["App Run", "CP4Apps", "CP4Systems", "DevOps", "ELM", "Project Harmony", "Runtimes", "Spectrum LSF", "UnifyBlue", "WAS", "WCA Java", "Web Hybrid ED"]
build_product_grid(slide, X_CENTER, AD_Y + Inches(0.35), ad_w, AD_H - Inches(0.35), ad_prods, cols=4, border_color=C_PURPLE_BORDER)
# App Int
create_box(slide, X_CENTER + ad_w + GAP, AD_Y, ad_w, Inches(0.3), "Application Integration", C_PURPLE_DARK, font_color=C_WHITE, font_size=11, bold=True)
ai_prods = ["API Connect", "APP Connect", "Aspera", "CP4I", "Connect:Direct", "DataPower", "DataPower Dashboard", "Event Automation", "FTM", "MQ", "Sterling B2B Integrator", "WebMethods"]
build_product_grid(slide, X_CENTER + ad_w + GAP, AD_Y + Inches(0.35), ad_w, AD_H - Inches(0.35), ai_prods, cols=4, border_color=C_PURPLE_BORDER)


# --- 5. RIGHT COLUMN ---
# IT Auto
IT_H = Inches(5.4)
create_box(slide, X_RIGHT, TOP_Y, RIGHT_W, Inches(0.35), "IT Automation & Finops", C_PURPLE_DARK, font_color=C_WHITE, font_size=10, bold=True)
it_prods = ["Ansible", "Apptio", "Cloud Pak for AIOps", "Cloudability", "Concert", "Flexera One", "HashiCorp Terraform", "Instana", "Kubecost", "Operations Insights", "Targetprocess", "Turbonomic", "Workload Automation"]
build_product_grid(slide, X_RIGHT, TOP_Y + Inches(0.4), RIGHT_W, IT_H - Inches(0.4), it_prods, cols=1, border_color=C_PURPLE_BORDER)

# Network
NET_Y = TOP_Y + IT_H + GAP
NET_H = (CA_H + GAP + P6_H + GAP + AD_H) - IT_H - GAP # Match Center Height
create_box(slide, X_RIGHT, NET_Y, RIGHT_W, Inches(0.35), "Network Mgmt", C_PURPLE_DARK, font_color=C_WHITE, font_size=11, bold=True)
net_prods = ["CP4NA", "Cloud Network Security", "Content Delivery Network", "Edge Application Manager", "HashiCorp Nomad", "Hybrid Cloud Mesh", "NS1 Connect", "SevOne"]
build_product_grid(slide, X_RIGHT, NET_Y + Inches(0.4), RIGHT_W, NET_H - Inches(0.4), net_prods, cols=1, border_color=C_PURPLE_BORDER)


# --- 6. RED HAT OPENSHIFT BANNER ---
# Spans Center + Right. White BG, Red Text, Dark Red Border.
RH_Y = AD_Y + AD_H + Inches(0.1)
RH_W = CENTER_W + GAP + RIGHT_W
create_box(slide, X_CENTER, RH_Y, RH_W, Inches(0.4), "Red Hat OpenShift", C_WHITE, outline_color=C_RED_DARK, outline_width=Pt(1.5), font_color=C_RED_DARK, bold=True, font_size=12)


# --- 7. INFRASTRUCTURE (Bottom Row) ---
INFRA_Y = RH_Y + Inches(0.5)
INFRA_H = Inches(1.5)

# A. Enterprise Storage (Aligned under Security)
create_box(slide, X_SEC, INFRA_Y, SEC_W, Inches(0.3), "Enterprise Storage", C_GREY_DARK, font_color=C_WHITE, font_size=9, bold=True)
es_prods = ["DS8000 Series", "SAN Directors", "Tape (Hydra & Jaguar)/VTS"]
build_product_grid(slide, X_SEC, INFRA_Y + Inches(0.35), SEC_W, INFRA_H - Inches(0.35), es_prods, cols=1, border_color=C_GREY_BORDER)

# B. Remaining Infrastructure Pillars (Aligned under Center + Right)
infra_remaining = [
    ("Data Resilience Storage", ["Scale", "Scale System", "Ceph", "CoS", "Defender/Protect", "Flash", "Fusion", "Fusion HCI", "Fusion HCI (on-prem)", "Hyperscaler", "SVC", "Ceph System", "Storage Insight", "Storage Virtualize", "Tape"]),
    ("Power", ["AIX", "IBM i", "Linux", "Oracle", "Red Hat OpenShift", "SAP"]),
    ("Z System", ["AI on Z", "IBM LinuxOne", "IBM zOS", "Z Monitoring Suite", "Z Security", "Z Software"]),
    ("Cloud", ["Cloud Financial Server", "Cloud Satellite", "Power Virtual Server", "Red Hat OpenShift", "SAP", "VMware"])
]

available_infra_w = CENTER_W + GAP + RIGHT_W
infra_unit_w = (available_infra_w - (GAP * 3)) / 9

cur_x = X_CENTER
for title, prods in infra_remaining:
    w_mult = 3 if "Resilience" in title else 2
    actual_w = (infra_unit_w * w_mult) + (GAP * (w_mult - 1))
    # Header
    create_box(slide, cur_x, INFRA_Y, actual_w, Inches(0.3), title, C_GREY_DARK, font_color=C_WHITE, font_size=9, bold=True)
    # Grid
    build_product_grid(slide, cur_x, INFRA_Y + Inches(0.35), actual_w, INFRA_H - Inches(0.35), prods, cols=w_mult, border_color=C_GREY_BORDER)
    cur_x += actual_w + GAP

# --- 8. FOOTER ---
# Light Grey BG, Dark Grey Border
FOOT_Y = INFRA_Y + INFRA_H + Inches(0.2)
FOOT_W = (SLIDE_WIDTH - (MARGIN_X * 2) - GAP) / 2
create_box(slide, MARGIN_X, FOOT_Y, FOOT_W, Inches(0.4), "IBM Technology Lifecycle Services (TLS)", C_GREY_LIGHT, outline_color=C_GREY_BORDER, outline_width=Pt(1.5), bold=True, font_size=11)
create_box(slide, MARGIN_X + FOOT_W + GAP, FOOT_Y, FOOT_W, Inches(0.4), "IBM Expert Labs (EL)", C_GREY_LIGHT, outline_color=C_GREY_BORDER, outline_width=Pt(1.5), bold=True, font_size=11)

prs.save("IBM_Product_Placemat.pptx")
print("Slide generated successfully.")
